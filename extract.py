import os
import re
import warnings
import sys
from typing import Optional, List

# 用于处理不同文件类型的库
try:
    import docx
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False
    warnings.warn("python-docx库未安装，将无法处理.docx文件")

try:
    import PyPDF2
    HAS_PDF = True
except ImportError:
    HAS_PDF = False
    warnings.warn("PyPDF2库未安装，将无法处理PDF文件")

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    warnings.warn("python-pptx库未安装，将无法处理.pptx文件")

try:
    import pandas as pd
    HAS_PANDAS = True
except ImportError:
    HAS_PANDAS = False
    warnings.warn("pandas库未安装，将无法处理Excel文件")

def extract_text_from_file(file_path: str, max_length: int = 10000) -> Optional[str]:
    """
    从文件中提取文本内容
    
    Args:
        file_path: 文件路径
        max_length: 最大提取文本长度，默认10000字符
    
    Returns:
        提取的文本内容，如果提取失败返回None
    """
    if not os.path.exists(file_path):
        print(f"文件不存在: {file_path}")
        return None
    
    file_extension = os.path.splitext(file_path)[1].lower()
    
    try:
        # 1. 文本文件 (.txt, .md, .csv等)
        if file_extension in ['.txt', '.md', '.csv', '.json', '.xml', '.html', '.htm']:
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    text = f.read(max_length * 2)  # 多读一些防止截断问题
            except UnicodeDecodeError:
                try:
                    with open(file_path, 'r', encoding='gbk') as f:
                        text = f.read(max_length * 2)
                except UnicodeDecodeError:
                    with open(file_path, 'r', encoding='latin-1') as f:
                        text = f.read(max_length * 2)
            
            # 如果是代码或标记语言，可以做一些基本清理
            if file_extension in ['.html', '.htm', '.xml']:
                text = re.sub(r'<[^>]+>', ' ', text)
                text = re.sub(r'\s+', ' ', text)
            
            return text[:max_length] if text else None
        
        # 2. Word文档 (.docx)
        elif file_extension == '.docx' and HAS_DOCX:
            doc = docx.Document(file_path)
            full_text = []
            for paragraph in doc.paragraphs:
                full_text.append(paragraph.text)
            text = '\n'.join(full_text)
            return text[:max_length] if text else None
        
        # 3. PDF文件 (.pdf)
        elif file_extension == '.pdf' and HAS_PDF:
            text = ''
            with open(file_path, 'rb') as f:
                pdf_reader = PyPDF2.PdfReader(f)
                for page in pdf_reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + '\n'
                    if len(text) >= max_length:
                        break
            return text[:max_length] if text else None
        
        # 4. PowerPoint文件 (.pptx)
        elif file_extension == '.pptx' and HAS_PPTX:
            prs = Presentation(file_path)
            full_text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        full_text.append(shape.text)
            text = '\n'.join(full_text)
            return text[:max_length] if text else None
        
        # 5. Excel文件 (.xlsx, .xls)
        elif file_extension in ['.xlsx', '.xls'] and HAS_PANDAS:
            try:
                # 读取所有sheet
                excel_file = pd.ExcelFile(file_path)
                all_text = []
                
                for sheet_name in excel_file.sheet_names:
                    df = pd.read_excel(file_path, sheet_name=sheet_name)
                    # 将DataFrame转换为文本
                    sheet_text = df.to_string(index=False)
                    all_text.append(f"--- Sheet: {sheet_name} ---\n{sheet_text}")
                    
                    if len('\n'.join(all_text)) >= max_length:
                        break
                
                text = '\n\n'.join(all_text)
                return text[:max_length] if text else None
            except Exception as e:
                print(f"读取Excel文件时出错: {e}")
                return None
        
        # 6. 其他文件类型尝试以二进制读取
        else:
            try:
                # 尝试以文本读取
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    text = f.read(max_length * 2)
                    return text[:max_length] if text else None
            except:
                # 如果是二进制文件，返回None表示没有文本
                return None
    
    except Exception as e:
        print(f"处理文件 {file_path} 时出错: {e}")
        return None


def save_file_text_to_txt(input_file_path: str, output_txt_path: str = None, max_length: int = 10000) -> bool:
    """
    提取文件中的文本并保存到txt文件
    
    Args:
        input_file_path: 输入文件路径
        output_txt_path: 输出txt文件路径（可选，默认与输入文件同名）
        max_length: 最大文本长度，默认10000字符
    
    Returns:
        bool: 是否成功保存
    """
    if not os.path.exists(input_file_path):
        print(f"错误: 文件 '{input_file_path}' 不存在")
        return False
    
    # 如果未指定输出路径，使用输入文件名
    if output_txt_path is None:
        base_name = os.path.splitext(os.path.basename(input_file_path))[0]
        output_txt_path = f"{base_name}_extracted.txt"
    
    # 提取文本
    extracted_text = extract_text_from_file(input_file_path, max_length)
    
    try:
        with open(output_txt_path, 'w', encoding='utf-8') as f:
            if extracted_text:
                # 写入提取的文本（截断到指定长度）
                truncated_text = extracted_text[:max_length]
                f.write(truncated_text)
                print(f"成功: 从 '{input_file_path}' 提取文本并保存到 '{output_txt_path}'")
                print(f"提取字符数: {len(truncated_text)}")
                if len(extracted_text) > max_length:
                    print(f"注意: 文本已截断，原始长度为 {len(extracted_text)} 字符")
            else:
                # 没有文本，只保存文件名
                file_name_only = os.path.basename(input_file_path)
                f.write(f"文件 '{file_name_only}' 中没有提取到文本内容。")
                print(f"注意: 文件 '{input_file_path}' 中没有提取到文本，只保存了文件名")
        return True
    
    except Exception as e:
        print(f"保存文件时出错: {e}")
        return False


def batch_extract_files(file_paths: list, output_dir: str = "extracted_texts", max_length: int = 10000) -> dict:
    """
    批量提取多个文件的文本
    
    Args:
        file_paths: 文件路径列表
        output_dir: 输出目录
        max_length: 最大文本长度
    
    Returns:
        dict: 处理结果统计
    """
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
    
    results = {
        'total': len(file_paths),
        'success': 0,
        'failed': 0,
        'no_text': 0,
        'failed_files': []
    }
    
    for file_path in file_paths:
        if os.path.exists(file_path):
            output_path = os.path.join(
                output_dir, 
                f"{os.path.splitext(os.path.basename(file_path))[0]}_extracted.txt"
            )
            
            success = save_file_text_to_txt(file_path, output_path, max_length)
            
            if success:
                # 检查是否提取到了文本
                try:
                    with open(output_path, 'r', encoding='utf-8') as f:
                        content = f.read()
                        if "没有提取到文本内容" in content:
                            results['no_text'] += 1
                        else:
                            results['success'] += 1
                except:
                    results['success'] += 1
            else:
                results['failed'] += 1
                results['failed_files'].append(file_path)
        else:
            print(f"文件不存在: {file_path}")
            results['failed'] += 1
            results['failed_files'].append(file_path)
    
    return results


def get_supported_extensions() -> List[str]:
    """获取支持的文件扩展名列表"""
    extensions = [
        '.txt', '.md', '.csv', '.json', '.xml', '.html', '.htm',
        '.py', '.js', '.java', '.cpp', '.c', '.php', '.rb'
    ]
    
    if HAS_DOCX:
        extensions.append('.docx')
    if HAS_PDF:
        extensions.append('.pdf')
    if HAS_PPTX:
        extensions.append('.pptx')
    if HAS_PANDAS:
        extensions.extend(['.xlsx', '.xls'])
    
    return extensions


def interactive_mode():
    """交互式模式"""
    print("=" * 60)
    print("文件文本提取工具 v2.0")
    print("支持的文件类型:", ", ".join(get_supported_extensions()))
    print("=" * 60)
    
    while True:
        print("\n" + "-" * 40)
        print("请选择操作:")
        print("1. 提取单个文件")
        print("2. 批量提取多个文件")
        print("3. 提取目录下所有支持的文件")
        print("4. 显示支持的文件类型")
        print("5. 检查依赖库")
        print("6. 退出")
        print("-" * 40)
        
        choice = input("请输入选项 (1-6): ").strip()
        
        if choice == "1":
            extract_single_file_interactive()
        elif choice == "2":
            extract_multiple_files_interactive()
        elif choice == "3":
            extract_directory_interactive()
        elif choice == "4":
            show_supported_extensions()
        elif choice == "5":
            check_dependencies()
        elif choice == "6":
            print("感谢使用，再见！")
            break
        else:
            print("无效选项，请重新输入！")


def extract_single_file_interactive():
    """交互式提取单个文件"""
    print("\n--- 提取单个文件 ---")
    
    while True:
        file_path = input("请输入文件路径 (或输入 'q' 返回主菜单): ").strip().strip('"').strip("'")
        
        if file_path.lower() == 'q':
            return
        
        if not os.path.exists(file_path):
            print(f"错误: 文件 '{file_path}' 不存在，请重新输入")
            continue
        
        # 获取最大长度
        max_length = get_max_length_input()
        
        # 获取输出路径
        default_output = f"{os.path.splitext(os.path.basename(file_path))[0]}_extracted.txt"
        output_path = input(f"请输入输出文件路径 (默认: {default_output}): ").strip()
        if not output_path:
            output_path = default_output
        
        print(f"\n开始处理...")
        print(f"输入文件: {file_path}")
        print(f"输出文件: {output_path}")
        print(f"最大长度: {max_length}")
        
        success = save_file_text_to_txt(file_path, output_path, max_length)
        
        if success:
            print("✓ 处理完成！")
        else:
            print("✗ 处理失败！")
        
        # 是否继续处理其他文件
        cont = input("\n是否继续处理其他文件？(y/n): ").strip().lower()
        if cont != 'y':
            break


def extract_multiple_files_interactive():
    """交互式提取多个文件"""
    print("\n--- 批量提取文件 ---")
    
    file_paths = []
    print("请输入文件路径列表 (每行一个，输入空行结束):")
    
    while True:
        file_path = input().strip().strip('"').strip("'")
        if not file_path:
            break
        
        if os.path.exists(file_path):
            file_paths.append(file_path)
        else:
            print(f"警告: 文件 '{file_path}' 不存在，已跳过")
    
    if not file_paths:
        print("未输入有效的文件路径")
        return
    
    print(f"\n共输入 {len(file_paths)} 个文件:")
    for i, fp in enumerate(file_paths, 1):
        print(f"{i:2d}. {fp}")
    
    # 获取最大长度
    max_length = get_max_length_input()
    
    # 获取输出目录
    default_output_dir = "batch_extracted"
    output_dir = input(f"请输入输出目录 (默认: {default_output_dir}): ").strip()
    if not output_dir:
        output_dir = default_output_dir
    
    confirm = input(f"\n确认开始处理 {len(file_paths)} 个文件？(y/n): ").strip().lower()
    if confirm != 'y':
        print("操作已取消")
        return
    
    print(f"\n开始批量处理...")
    results = batch_extract_files(file_paths, output_dir, max_length)
    
    print("\n" + "=" * 40)
    print("批量处理完成！")
    print(f"总计处理: {results['total']} 个文件")
    print(f"成功提取文本: {results['success']} 个")
    print(f"无文本内容: {results['no_text']} 个")
    print(f"处理失败: {results['failed']} 个")
    
    if results['failed_files']:
        print(f"\n失败的文件列表:")
        for failed_file in results['failed_files']:
            print(f"  - {failed_file}")


def extract_directory_interactive():
    """交互式提取目录下所有文件"""
    print("\n--- 提取目录下所有支持的文件 ---")
    
    while True:
        dir_path = input("请输入目录路径 (或输入 '.' 使用当前目录): ").strip()
        if dir_path == '.' or not dir_path:
            dir_path = os.getcwd()
        
        if not os.path.isdir(dir_path):
            print(f"错误: '{dir_path}' 不是有效目录")
            continue
        
        # 获取最大长度
        max_length = get_max_length_input()
        
        # 查找所有支持的文件
        supported_ext = get_supported_extensions()
        all_files = []
        
        print(f"\n正在扫描目录: {dir_path}")
        for root, dirs, files in os.walk(dir_path):
            for file in files:
                file_ext = os.path.splitext(file)[1].lower()
                if file_ext in supported_ext:
                    file_path = os.path.join(root, file)
                    all_files.append(file_path)
        
        if not all_files:
            print(f"目录 '{dir_path}' 中未找到支持的文件")
            return
        
        print(f"找到 {len(all_files)} 个支持的文件")
        
        # 显示前20个文件
        print("\n前20个文件:")
        for i, file_path in enumerate(all_files[:20], 1):
            print(f"{i:2d}. {os.path.relpath(file_path, dir_path)}")
        
        if len(all_files) > 20:
            print(f"... 还有 {len(all_files) - 20} 个文件")
        
        # 获取输出目录
        default_output_dir = os.path.join(dir_path, "extracted_texts")
        output_dir = input(f"\n请输入输出目录 (默认: {default_output_dir}): ").strip()
        if not output_dir:
            output_dir = default_output_dir
        
        confirm = input(f"\n确认处理这 {len(all_files)} 个文件？(y/n): ").strip().lower()
        if confirm != 'y':
            print("操作已取消")
            return
        
        print(f"\n开始处理...")
        results = batch_extract_files(all_files, output_dir, max_length)
        
        print("\n" + "=" * 40)
        print("处理完成！")
        print(f"目录: {dir_path}")
        print(f"总计: {results['total']} 个文件")
        print(f"成功: {results['success']} 个")
        print(f"无文本: {results['no_text']} 个")
        print(f"失败: {results['failed']} 个")
        
        # 保存处理摘要
        summary_file = os.path.join(output_dir, "处理摘要.txt")
        try:
            with open(summary_file, 'w', encoding='utf-8') as f:
                f.write("文件文本提取处理摘要\n")
                f.write("=" * 50 + "\n")
                f.write(f"处理目录: {dir_path}\n")
                f.write(f"输出目录: {output_dir}\n")
                f.write(f"最大文本长度: {max_length}\n")
                f.write(f"处理时间: {get_current_time()}\n")
                f.write(f"处理文件总数: {results['total']}\n")
                f.write(f"成功提取文本: {results['success']}\n")
                f.write(f"无文本内容: {results['no_text']}\n")
                f.write(f"处理失败: {results['failed']}\n")
                
                if results['failed_files']:
                    f.write("\n失败的文件:\n")
                    for failed_file in results['failed_files']:
                        f.write(f"- {failed_file}\n")
            
            print(f"处理摘要已保存到: {summary_file}")
        except Exception as e:
            print(f"保存处理摘要时出错: {e}")
        
        break


def get_max_length_input() -> int:
    """获取用户输入的最大长度"""
    while True:
        max_length_input = input("请输入最大文本长度 (默认10000): ").strip()
        
        if not max_length_input:
            return 10000
        
        try:
            max_length = int(max_length_input)
            if max_length <= 0:
                print("最大长度必须大于0，请重新输入")
                continue
            return max_length
        except ValueError:
            print("请输入有效的数字，请重新输入")


def show_supported_extensions():
    """显示支持的文件扩展名"""
    print("\n支持的文件类型:")
    print("-" * 30)
    
    extensions = get_supported_extensions()
    extensions.sort()
    
    for i, ext in enumerate(extensions, 1):
        description = {
            '.txt': '文本文件',
            '.md': 'Markdown文件',
            '.csv': 'CSV文件',
            '.json': 'JSON文件',
            '.xml': 'XML文件',
            '.html': 'HTML文件',
            '.htm': 'HTML文件',
            '.docx': 'Word文档',
            '.pdf': 'PDF文件',
            '.pptx': 'PowerPoint文件',
            '.xlsx': 'Excel文件',
            '.xls': 'Excel文件',
            '.py': 'Python脚本',
            '.js': 'JavaScript文件',
            '.java': 'Java文件',
            '.cpp': 'C++文件',
            '.c': 'C文件',
            '.php': 'PHP文件',
            '.rb': 'Ruby文件'
        }.get(ext, '其他文件')
        
        # 检查是否已安装对应的库
        available = "✓"
        if ext == '.docx' and not HAS_DOCX:
            available = "✗ (需要安装 python-docx)"
        elif ext == '.pdf' and not HAS_PDF:
            available = "✗ (需要安装 PyPDF2)"
        elif ext == '.pptx' and not HAS_PPTX:
            available = "✗ (需要安装 python-pptx)"
        elif ext in ['.xlsx', '.xls'] and not HAS_PANDAS:
            available = "✗ (需要安装 pandas)"
        
        print(f"{ext:8} - {description:15} {available}")
    
    print("-" * 30)


def check_dependencies():
    """检查依赖库"""
    print("\n依赖库检查:")
    print("-" * 30)
    
    dependencies = [
        ("python-docx", HAS_DOCX, "用于处理Word文档 (.docx)"),
        ("PyPDF2", HAS_PDF, "用于处理PDF文件 (.pdf)"),
        ("python-pptx", HAS_PPTX, "用于处理PowerPoint文件 (.pptx)"),
        ("pandas", HAS_PANDAS, "用于处理Excel文件 (.xlsx, .xls)")
    ]
    
    all_installed = True
    for name, installed, description in dependencies:
        status = "✓ 已安装" if installed else "✗ 未安装"
        print(f"{name:15} - {status:15} {description}")
        if not installed:
            all_installed = False
    
    print("-" * 30)
    
    if not all_installed:
        print("\n要安装缺失的依赖库，可以使用以下命令:")
        print("pip install python-docx PyPDF2 python-pptx pandas")
    
    return all_installed


def get_current_time() -> str:
    """获取当前时间字符串"""
    from datetime import datetime
    return datetime.now().strftime("%Y-%m-%d %H:%M:%S")


def command_line_mode():
    """命令行模式"""
    if len(sys.argv) < 2:
        print("用法:")
        print("  python script.py <文件路径> [输出路径] [最大长度]")
        print("  python script.py --batch <文件1> <文件2> ...")
        print("  python script.py --dir <目录路径>")
        print("  python script.py --interactive")
        print("\n示例:")
        print("  python script.py document.docx")
        print("  python script.py input.pdf output.txt 5000")
        print("  python script.py --batch file1.txt file2.docx")
        print("  python script.py --dir ./documents")
        return
    
    if sys.argv[1] == "--interactive" or sys.argv[1] == "-i":
        interactive_mode()
    elif sys.argv[1] == "--batch" or sys.argv[1] == "-b":
        if len(sys.argv) < 3:
            print("错误: 请指定要处理的文件")
            return
        file_paths = sys.argv[2:]
        max_length = 10000
        if len(sys.argv) > 3 and sys.argv[-1].isdigit():
            max_length = int(sys.argv[-1])
            file_paths = file_paths[:-1]
        results = batch_extract_files(file_paths, "batch_extracted", max_length)
        print(f"处理完成: {results['success']}/{results['total']} 成功")
    elif sys.argv[1] == "--dir" or sys.argv[1] == "-d":
        if len(sys.argv) < 3:
            print("错误: 请指定目录路径")
            return
        dir_path = sys.argv[2]
        max_length = 10000
        if len(sys.argv) > 3 and sys.argv[-1].isdigit():
            max_length = int(sys.argv[-1])
        
        supported_ext = get_supported_extensions()
        all_files = []
        for root, dirs, files in os.walk(dir_path):
            for file in files:
                file_ext = os.path.splitext(file)[1].lower()
                if file_ext in supported_ext:
                    file_path = os.path.join(root, file)
                    all_files.append(file_path)
        
        if all_files:
            results = batch_extract_files(all_files, "extracted_texts", max_length)
            print(f"处理完成: 找到 {len(all_files)} 个文件, {results['success']} 成功")
        else:
            print(f"目录 '{dir_path}' 中未找到支持的文件")
    else:
        # 单个文件模式
        file_path = sys.argv[1]
        output_path = sys.argv[2] if len(sys.argv) > 2 else None
        max_length = int(sys.argv[3]) if len(sys.argv) > 3 and sys.argv[3].isdigit() else 10000
        
        if not os.path.exists(file_path):
            print(f"错误: 文件 '{file_path}' 不存在")
            return
        
        success = save_file_text_to_txt(file_path, output_path, max_length)
        if success:
            print("处理成功！")
        else:
            print("处理失败！")


if __name__ == "__main__":
    # 支持两种使用方式：命令行参数或交互式
    if len(sys.argv) > 1:
        command_line_mode()
    else:
        interactive_mode()
